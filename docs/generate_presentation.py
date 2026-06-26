"""
Generates the final Actu'AI presentation (PPTX).

Run with:
    uv run --with python-pptx python docs/generate_presentation.py

Visual language mirrors the team's earlier draft deck (numbered section
badge, navy serif titles, white cards, footer "Actu'AI - ECE Paris") but the
content is rewritten for the COMPLETED project: missions are described as
delivered and validated, with concrete evidence from the end-to-end Docker
Compose integration test pass.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---------------------------------------------------------------- palette --
NAVY = RGBColor(0x14, 0x28, 0x4B)
NAVY_DARK = RGBColor(0x0B, 0x1B, 0x3A)
BLUE = RGBColor(0x2F, 0x6F, 0xED)
BLUE_LIGHT = RGBColor(0xEA, 0xF1, 0xFD)
ORANGE = RGBColor(0xE8, 0x7B, 0x1E)
ORANGE_LIGHT = RGBColor(0xFD, 0xEE, 0xDC)
GREEN = RGBColor(0x1A, 0x7A, 0x3E)
GREEN_LIGHT = RGBColor(0xE3, 0xF5, 0xE9)
PURPLE = RGBColor(0x7C, 0x5C, 0xD4)
PURPLE_LIGHT = RGBColor(0xEF, 0xEA, 0xFB)
TEAL = RGBColor(0x16, 0x8A, 0x7C)
TEAL_LIGHT = RGBColor(0xE1, 0xF5, 0xF2)
RED = RGBColor(0xD1, 0x3B, 0x3B)
RED_LIGHT = RGBColor(0xFC, 0xE8, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF4, 0xF6, 0xFB)
GRAY_TXT = RGBColor(0x4B, 0x55, 0x66)
GRAY_LINE = RGBColor(0xDD, 0xE3, 0xEC)

FONT_HEAD = "Georgia"
FONT_BODY = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06):
    sp = slide.shapes.add_shape(shape_type, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    no_line(sp)
    sp.shadow.inherit = False
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def soft_shadow(shape):
    el = shape._element.spPr
    effectLst = el.makeelement(qn("a:effectLst"), {})
    outer = el.makeelement(qn("a:outerShdw"), {
        "blurRad": "90000", "dist": "30000", "dir": "5400000", "rotWithShape": "0",
    })
    clr = el.makeelement(qn("a:srgbClr"), {"val": "1B2A4A"})
    alpha = el.makeelement(qn("a:alpha"), {"val": "18000"})
    clr.append(alpha)
    outer.append(clr)
    effectLst.append(outer)
    el.append(effectLst)


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def run_text(p, text, size, color, bold=False, italic=False, font=FONT_BODY, spacing=None):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    if spacing is not None:
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(spacing))
    return r


def simple_text(slide, x, y, w, h, text, size, color, bold=False, italic=False,
                 align=PP_ALIGN.LEFT, font=FONT_BODY, anchor=MSO_ANCHOR.TOP, spacing=None,
                 line_spacing=None):
    tb, tf = textbox(slide, x, y, w, h, anchor=anchor)
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run_text(p, text, size, color, bold, italic, font, spacing)
    return tb


def bullets(slide, x, y, w, h, items, size=12.5, color=GRAY_TXT, gap=6, bold_lead=None,
            line_spacing=1.12, bullet_color=None):
    tb, tf = textbox(slide, x, y, w, h)
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.line_spacing = line_spacing
        run_text(p, "•  ", size, bullet_color or BLUE, bold=True, font=FONT_BODY)
        if isinstance(item, tuple):
            lead, rest = item
            run_text(p, lead, size, NAVY, bold=True, font=FONT_BODY)
            run_text(p, rest, size, color, font=FONT_BODY)
        else:
            run_text(p, item, size, color, font=FONT_BODY)
    return tb


def footer(slide, page_no=None):
    simple_text(slide, Inches(8.6), Inches(7.12), Inches(4.4), Inches(0.3),
                "Actu'AI · ECE Paris — Final Delivery", 9, RGBColor(0x9A, 0xA3, 0xB2),
                align=PP_ALIGN.RIGHT, italic=True)


def header(slide, number, eyebrow, title, title_size=27, title_color=NAVY, title_w=Inches(10.8)):
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), Inches(0.42), Inches(0.62), Inches(0.62))
    circ.fill.solid()
    circ.fill.fore_color.rgb = BLUE
    no_line(circ)
    circ.shadow.inherit = False
    tf = circ.text_frame
    tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run_text(p, number, 16, WHITE, bold=True, font=FONT_BODY)

    simple_text(slide, Inches(1.35), Inches(0.40), Inches(8), Inches(0.3),
                eyebrow, 11, BLUE, bold=True, spacing=150)
    simple_text(slide, Inches(1.33), Inches(0.68), title_w, Inches(0.6),
                title, title_size, title_color, bold=True, font=FONT_HEAD)


def card(slide, x, y, w, h, fill=WHITE, line=GRAY_LINE):
    sp = rect(slide, x, y, w, h, fill, radius=0.045)
    sp.line.color.rgb = line
    sp.line.width = Pt(0.75)
    soft_shadow(sp)
    return sp


def icon_badge(slide, x, y, d, bg, glyph=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    sp.fill.solid()
    sp.fill.fore_color.rgb = bg
    no_line(sp)
    sp.shadow.inherit = False
    return sp


def gradient_bg(slide, c1, c2, angle=45):
    slide.background.fill.gradient()
    stops = slide.background.fill.gradient_stops
    stops[0].color.rgb = c1
    stops[0].position = 0.0
    stops[1].color.rgb = c2
    stops[1].position = 1.0
    slide.background.fill.gradient_angle = angle


# ============================================================ SLIDE 1 =====
s = add_slide()
gradient_bg(s, NAVY_DARK, BLUE)
logo = card(s, Inches(0.6), Inches(0.5), Inches(2.0), Inches(0.85), fill=WHITE, line=WHITE)
simple_text(s, Inches(0.6), Inches(0.5), Inches(2.0), Inches(0.85), "ECE", 22, NAVY, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)

simple_text(s, Inches(0.9), Inches(2.35), Inches(11), Inches(0.4),
            "AEROSPACE  ·  MULTI-AGENT AI  ·  DATA SOVEREIGNTY  ·  DELIVERED", 13,
            RGBColor(0xCF, 0xDE, 0xF7), bold=True, spacing=150)
simple_text(s, Inches(0.85), Inches(2.7), Inches(10), Inches(1.3),
            "Actu'AI", 64, WHITE, bold=True, font=FONT_HEAD)
simple_text(s, Inches(0.9), Inches(3.95), Inches(9.6), Inches(1.0),
            "Final delivery — a multi-agent AI system that automates the non-value-added "
            "administrative work of an A350 thrust-reverser actuation supply chain, validated "
            "end-to-end on its own Docker Compose stack, while keeping all sensitive aerospace "
            "data on local infrastructure.",
            14.5, RGBColor(0xE4, 0xEC, 0xFB), line_spacing=1.2)

infobox = card(s, Inches(0.9), Inches(5.55), Inches(7.6), Inches(1.25), fill=WHITE, line=WHITE)
tb, tf = textbox(s, Inches(1.2), Inches(5.72), Inches(7.1), Inches(0.95))
p = tf.paragraphs[0]
run_text(p, "TEAM   ", 10.5, BLUE, bold=True, spacing=100)
run_text(p, "Jarfino HOUNGBADJI  ·  SOW Achta Demba  ·  Fatoumata IBRAHIM BELKO", 11.5, NAVY)
p2 = tf.add_paragraph(); p2.space_before = Pt(8)
run_text(p2, "SUPERVISOR   ", 10.5, ORANGE, bold=True, spacing=100)
run_text(p2, "GUETARI Skander          ", 11.5, NAVY)
run_text(p2, "STATUS   ", 10.5, GREEN, bold=True, spacing=100)
run_text(p2, "Completed & integration-tested", 11.5, NAVY, bold=True)

simple_text(s, Inches(8.9), Inches(6.95), Inches(3.8), Inches(0.3),
            "Source code  github.com/jaja07/ActuAI", 10.5, RGBColor(0xCF, 0xDE, 0xF7),
            align=PP_ALIGN.RIGHT, bold=True)

# ============================================================ SLIDE 2 =====
s = add_slide(); set_bg(s, BG)
header(s, "02", "EXECUTIVE SUMMARY", "Final delivery at a glance")

c = card(s, Inches(0.55), Inches(1.35), Inches(6.45), Inches(5.4))
simple_text(s, Inches(0.85), Inches(1.62), Inches(5.9), Inches(0.4), "What was delivered", 15, NAVY, bold=True)
tb, tf = textbox(s, Inches(0.85), Inches(2.1), Inches(5.85), Inches(2.4))
p = tf.paragraphs[0]; p.line_spacing = 1.2
run_text(p, "Actu'AI now runs as a complete multi-agent system that ingests supplier "
             "emails and audit requests, drafts the required SAP updates, alerts and documents "
             "automatically, and surfaces every action to a human expert for validation. ", 12.5, GRAY_TXT)
p2 = tf.add_paragraph(); p2.space_before = Pt(10); p2.line_spacing = 1.2
run_text(p2, "All five operational missions are implemented, wired end-to-end, and were "
              "exercised against a live Docker Compose stack — real PostgreSQL, real Qdrant, a "
              "simulated SAP ERP, and the React validation dashboard — not just unit tests.", 12.5, GRAY_TXT)

icons_y = Inches(4.75)
labels = ["Eliminates NVA\ntasks", "Local data\nsovereignty", "Human-in-the-\nloop", "Validated end-\nto-end"]
colors = [GREEN, ORANGE, BLUE, PURPLE]
for i, (lab, col) in enumerate(zip(labels, colors)):
    bx = Inches(0.85 + i * 1.5)
    icon_badge(s, bx, icons_y, Inches(0.55), {GREEN: GREEN_LIGHT, ORANGE: ORANGE_LIGHT, BLUE: BLUE_LIGHT, PURPLE: PURPLE_LIGHT}[col])
    simple_text(s, bx, icons_y, Inches(0.55), Inches(0.55), "✓", 16, col, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    simple_text(s, bx - Inches(0.3), icons_y + Inches(0.65), Inches(1.15), Inches(0.6),
                lab, 9.5, GRAY_TXT, align=PP_ALIGN.CENTER, line_spacing=1.05)

stat_x = Inches(7.25)
stats = [
    ("30–60", "supplier emails handled per employee, per day", BLUE),
    ("3 h 10", "NVA delay on a single delay alert — now near-instant", ORANGE),
    ("5 / 5", "core actuation missions implemented and validated", GREEN),
    ("0", "regressions across 11 automated backend tests", PURPLE),
]
sh = Inches(1.22)
for i, (num, lab, col) in enumerate(stats):
    cy = Inches(1.35) + i * (sh + Inches(0.08))
    cc = card(s, stat_x, cy, Inches(5.55), sh)
    icon_badge(s, stat_x + Inches(0.25), cy + Inches(0.3), Inches(0.62), {BLUE: BLUE_LIGHT, ORANGE: ORANGE_LIGHT, GREEN: GREEN_LIGHT, PURPLE: PURPLE_LIGHT}[col])
    simple_text(s, stat_x + Inches(1.1), cy + Inches(0.12), Inches(1.6), Inches(0.55), num, 24, col, bold=True, font=FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE)
    simple_text(s, stat_x + Inches(2.75), cy + Inches(0.18), Inches(2.65), Inches(0.85), lab, 11, GRAY_TXT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

footer(s)

# ============================================================ SLIDE 3 =====
s = add_slide(); set_bg(s, BG)
header(s, "03", "THE DOMAIN", "Aerospace actuation & the A350 elecTRAS")

c = card(s, Inches(0.55), Inches(1.35), Inches(7.55), Inches(2.55))
simple_text(s, Inches(0.85), Inches(1.58), Inches(7), Inches(0.4), "From control signal to physical motion", 14, NAVY, bold=True)
bullets(s, Inches(0.85), Inches(2.05), Inches(6.95), Inches(1.75), [
    "If the onboard computer is a machine's brain, the actuation system is its "
    "muscles — converting a control signal into mechanical action.",
    "This project targets the elecTRAS, the Thrust Reverser Actuation System of "
    "the A350, which deploys the engine cowlings on landing to decelerate the aircraft.",
], size=12, gap=8)

c2 = card(s, Inches(0.55), Inches(4.05), Inches(7.55), Inches(2.7))
simple_text(s, Inches(0.85), Inches(4.28), Inches(7), Inches(0.4), "Precision components managed by the service", 14, NAVY, bold=True)
bullets(s, Inches(0.85), Inches(4.75), Inches(6.95), Inches(1.85), [
    ("Electromechanical actuators — ", "convert electrical commands into the force that moves the cowlings."),
    ("Gearmotors — ", "geared drive units delivering controlled torque for deployment."),
    ("Embedded electronic boards — ", "on-board control electronics for actuation logic and feedback."),
    ("Precision mechanical parts — ", "high-tolerance machined parts integrated into the nacelle."),
], size=11.3, gap=5)

c3 = card(s, Inches(8.3), Inches(1.35), Inches(4.45), Inches(5.4), fill=NAVY, line=NAVY)
simple_text(s, Inches(8.6), Inches(1.6), Inches(4), Inches(0.4), "Why it is critical", 15, WHITE, bold=True)
tb, tf = textbox(s, Inches(8.6), Inches(2.15), Inches(3.9), Inches(1.1))
p = tf.paragraphs[0]; p.line_spacing = 1.2
run_text(p, "Handling these parts must strictly follow rigorous aerospace quality and "
             "safety regulations. Every operation is governed by:", 11.5, RGBColor(0xCF, 0xDE, 0xF7))

for i, (tag, desc) in enumerate([
    ("EN 9100", "European quality management standard for aviation, space and defence."),
    ("AS 9100", "Equivalent aerospace standard governing manufacturing and traceability."),
]):
    by = Inches(3.45) + i * Inches(1.05)
    bb = rect(s, Inches(8.6), by, Inches(3.9), Inches(0.9), RGBColor(0x1F, 0x3D, 0x70), radius=0.12)
    simple_text(s, Inches(8.8), by + Inches(0.08), Inches(3.5), Inches(0.3), tag, 13, WHITE, bold=True)
    simple_text(s, Inches(8.8), by + Inches(0.38), Inches(3.5), Inches(0.5), desc, 10, RGBColor(0xCF, 0xDE, 0xF7), line_spacing=1.05)

simple_text(s, Inches(8.6), Inches(5.75), Inches(3.9), Inches(0.85),
            "Components feed directly into the Airbus production schedule — any slip "
            "cascades downstream into an AOG event.", 10.5, RGBColor(0xCF, 0xDE, 0xF7), italic=True, line_spacing=1.15)
footer(s)

# ============================================================ SLIDE 4 =====
s = add_slide(); set_bg(s, BG)
header(s, "04", "THE CORE PROBLEM", "The ERP–communication gap")

c1 = card(s, Inches(0.55), Inches(1.35), Inches(11.2), Inches(1.85))
icon_badge(s, Inches(0.95), Inches(1.6), Inches(0.55), BLUE_LIGHT)
simple_text(s, Inches(0.95), Inches(1.6), Inches(0.55), Inches(0.55), "🗄", 16, BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
simple_text(s, Inches(1.7), Inches(1.6), Inches(3.6), Inches(0.35), "Rigid structured world", 12.5, NAVY, bold=True)
simple_text(s, Inches(1.7), Inches(2.0), Inches(3.6), Inches(1.0), "SAP ERP: forecast dates, part references, "
            "PO numbers, serial numbers — the system of record.", 11, GRAY_TXT, line_spacing=1.15)

simple_text(s, Inches(5.5), Inches(1.55), Inches(2.3), Inches(0.4), "MANUAL", 12, ORANGE, bold=True, align=PP_ALIGN.CENTER)
simple_text(s, Inches(5.5), Inches(1.9), Inches(2.3), Inches(0.4), "⇄  BRIDGE = NVA", 13, ORANGE, bold=True, align=PP_ALIGN.CENTER)
simple_text(s, Inches(5.5), Inches(2.4), Inches(2.3), Inches(0.6), "→ now automated by\nthe Transactional &\nInvestigative agents", 9.5, GRAY_TXT, align=PP_ALIGN.CENTER, italic=True, line_spacing=1.05)

icon_badge(s, Inches(8.4), Inches(1.6), Inches(0.55), ORANGE_LIGHT)
simple_text(s, Inches(8.4), Inches(1.6), Inches(0.55), Inches(0.55), "≡", 16, ORANGE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
simple_text(s, Inches(9.15), Inches(1.6), Inches(2.5), Inches(0.35), "Unstructured world", 12.5, NAVY, bold=True)
simple_text(s, Inches(9.15), Inches(2.0), Inches(2.5), Inches(1.0), "Daily emails, PDF documents, Excel "
            "files, scanned paper — the real flow of information.", 11, GRAY_TXT, line_spacing=1.15)

simple_text(s, Inches(0.55), Inches(3.5), Inches(11), Inches(0.4),
            "Concrete example — a single 8-day delay alert on a critical A350 component (before automation)",
            13, NAVY, bold=True)

steps = [("1", "07:50", "Supplier emails to report\nan 8-day delay", BLUE),
         ("2", "09:20", "Employee reads\nthe email", ORANGE),
         ("3", "09:40", "Data is manually entered\ninto SAP", ORANGE),
         ("4", "10:00", "Planning manager\nis alerted", ORANGE),
         ("5", "11:00", "Airbus line learns of it\nin a meeting", GREEN)]
sw = Inches(2.16)
for i, (n, t, lab, col) in enumerate(steps):
    sx = Inches(0.55) + i * sw
    circ = slide_circ = s.shapes.add_shape(MSO_SHAPE.OVAL, sx + Inches(0.78), Inches(4.05), Inches(0.55), Inches(0.55))
    circ.fill.solid(); circ.fill.fore_color.rgb = col; no_line(circ); circ.shadow.inherit = False
    tf = circ.text_frame; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run_text(p, n, 14, WHITE, bold=True)
    simple_text(s, sx, Inches(4.7), sw, Inches(0.35), t, 14, col, bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD)
    simple_text(s, sx + Inches(0.1), Inches(5.05), sw - Inches(0.2), Inches(0.6), lab, 10, GRAY_TXT, align=PP_ALIGN.CENTER, line_spacing=1.1)
    if i < 4:
        ln = s.shapes.add_connector(1, sx + Inches(1.9), Inches(4.32), sx + sw + Inches(0.78), Inches(4.32))
        ln.line.color.rgb = GRAY_LINE
        ln.line.width = Pt(1.5)

res = rect(s, Inches(0.55), Inches(5.95), Inches(11.2), Inches(0.75), NAVY, radius=0.1)
simple_text(s, Inches(0.9), Inches(5.95), Inches(10.6), Inches(0.75),
            "Result (then): a 3 h 10 processing delay that was entirely non-value-added. "
            "Result (now): drafted, validated and acted on in minutes — see Mission 1, next.",
            12, WHITE, bold=False, anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ============================================================ SLIDE 5 =====
s = add_slide(); set_bg(s, BG)
header(s, "05", "BUSINESS PROCESS LAYER", "Five missions — now fully automated")
simple_text(s, Inches(0.55), Inches(1.3), Inches(11.5), Inches(0.5),
            "The Actuation service runs five core missions every day. Each one used to generate a "
            "stream of non-value-added work; each is now handled end-to-end by Actu'AI's agents.",
            12, GRAY_TXT, line_spacing=1.15)

missions = [
    ("1", "Supply Chain\nMonitoring", "Track deliveries against contractual dates.", "M1", BLUE, BLUE_LIGHT),
    ("2", "Production Schedule\nCoordination", "Synchronise deliveries with the Airbus build plan.", "M2", ORANGE, ORANGE_LIGHT),
    ("3", "Quality &\nNon-Conformity", "Document defects under EN9100/AS9100.", "M3", GREEN, GREEN_LIGHT),
    ("4", "Technical Documentation\nControl", "Manage regulatory records & certificates.", "M4", PURPLE, PURPLE_LIGHT),
    ("5", "End-to-End\nTraceability", "Reconstruct a part's full history, order to nacelle.", "M5", TEAL, TEAL_LIGHT),
]
mw = Inches(2.18)
for i, (n, title, desc, tag, col, lcol) in enumerate(missions):
    mx = Inches(0.55) + i * (mw + Inches(0.06))
    cc = card(s, mx, Inches(2.0), mw, Inches(3.4))
    icon_badge(s, mx + Inches(0.18), Inches(2.25), Inches(0.5), lcol)
    simple_text(s, mx + Inches(0.18), Inches(2.25), Inches(0.5), Inches(0.5), n, 15, col, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    simple_text(s, mx + Inches(0.18), Inches(2.9), mw - Inches(0.35), Inches(0.7), title, 12.5, NAVY, bold=True, line_spacing=1.05)
    simple_text(s, mx + Inches(0.18), Inches(3.65), mw - Inches(0.35), Inches(0.9), desc, 10, GRAY_TXT, line_spacing=1.12)
    tagb = rect(s, mx + Inches(0.18), Inches(4.7), Inches(0.9), Inches(0.35), lcol, radius=0.3)
    simple_text(s, mx + Inches(0.18), Inches(4.7), Inches(0.9), Inches(0.35), tag, 10.5, col, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    simple_text(s, mx + Inches(0.18), Inches(5.2), mw - Inches(0.35), Inches(0.35), "✓ Delivered", 10.5, GREEN, bold=True)

closing = rect(s, Inches(0.55), Inches(5.8), Inches(11.2), Inches(0.95), BLUE, radius=0.08)
simple_text(s, Inches(0.9), Inches(5.8), Inches(10.6), Inches(0.95),
            "Every manual step that used to bridge SAP and the daily document flow is now an "
            "automated draft, paused for one human validation click — across all five missions.",
            12.5, WHITE, italic=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
footer(s)


def mission_deep_dive(num, eyebrow, title, left, right):
    """left/right: dict(title, color, light, bullets:list, proof:str)"""
    s = add_slide(); set_bg(s, BG)
    header(s, num, eyebrow, title, title_size=25)
    panels = [left, right] if right else [left]
    pw = Inches(11.2) if not right else Inches(5.55)
    for i, panel in enumerate(panels):
        px = Inches(0.55) + i * (pw + Inches(0.1))
        cc = card(s, px, Inches(1.35), pw, Inches(5.45))
        icon_badge(s, px + Inches(0.3), Inches(1.6), Inches(0.55), panel["light"])
        simple_text(s, px + Inches(1.0), Inches(1.62), pw - Inches(1.3), Inches(0.55),
                    panel["title"], 14.5, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        simple_text(s, px + Inches(0.3), Inches(2.3), pw - Inches(0.6), Inches(0.3),
                    "Delivered capability", 10.5, panel["color"], bold=True, spacing=80)
        bullets(s, px + Inches(0.3), Inches(2.65), pw - Inches(0.6), Inches(2.55), panel["bullets"],
                size=11.3, gap=7, bullet_color=panel["color"])
        proofbox = rect(s, px + Inches(0.3), Inches(5.35), pw - Inches(0.6), Inches(1.3), panel["light"], radius=0.1)
        simple_text(s, px + Inches(0.5), Inches(5.45), pw - Inches(1.0), Inches(0.3),
                    "Validated in integration test", 10, panel["color"], bold=True)
        simple_text(s, px + Inches(0.5), Inches(5.72), pw - Inches(1.0), Inches(0.85),
                    panel["proof"], 10.3, GRAY_TXT, line_spacing=1.12)
    footer(s)
    return s


# ============================================================ SLIDE 6 =====
mission_deep_dive(
    "06", "MISSION 1 · DELIVERED", "Supply Chain Monitoring — automatic SAP sync",
    {
        "title": "Transactional Agent", "color": BLUE, "light": BLUE_LIGHT,
        "bullets": [
            "Reads each incoming supplier email and classifies it without manual sorting "
            "of the daily 30–60 message flood.",
            "Extracts the delivery date and status, and drafts a structured SAP update "
            "payload — eliminating manual double data entry.",
            "Pauses for explicit human validation before anything is written to SAP.",
        ],
        "proof": "A simulated delay email for a real PO was ingested through the live API; "
                 "the agent drafted a SAP_UPDATE task, and on approval the mock SAP's "
                 "expected_delivery_date was confirmed updated.",
    },
    None,
)

# ============================================================ SLIDE 7 =====
mission_deep_dive(
    "07", "MISSIONS 2 & 3 · DELIVERED", "Production Scheduling & Quality Management",
    {
        "title": "M2 · AOG Risk Detection", "color": ORANGE, "light": ORANGE_LIGHT,
        "bullets": [
            "Cross-references every delivery update against the Airbus production "
            "schedule for the same part reference.",
            "Flags a delay against the assembly drop-dead date automatically, generating "
            "a Red AOG alert — no more D-Day surprises.",
            "On approval, triggers an escalation email to the supplier/transporter "
            "requesting expedited shipping.",
        ],
        "proof": "The same delay event was cross-checked against a real production "
                 "schedule row; the AOG_ALERT task was created in parallel, and on "
                 "approval a real escalation email landed in the outbox with the "
                 "correct drop-dead delta.",
    },
    {
        "title": "M3 · Quality & FNC", "color": GREEN, "light": GREEN_LIGHT,
        "bullets": [
            "Detects a non-conformity request and pre-fills the Non-Conformance Report "
            "(part reference, supplier, PO, reception date) straight from SAP data.",
            "Eliminates manual re-typing across quality registers and supplier reports.",
            "On approval, posts the FNC directly to the SAP quality-notification endpoint.",
        ],
        "proof": "\"Créer FNC pour rayure...\" on a real PO produced a pre-filled FNC "
                 "draft; approving it created a real Quality Notification record "
                 "in the mock SAP, confirmed via the BAPI.",
    },
)

# ============================================================ SLIDE 8 =====
mission_deep_dive(
    "08", "MISSIONS 4 & 5 · DELIVERED", "Documentation Search & End-to-End Traceability",
    {
        "title": "M4 · Investigative Agent", "color": PURPLE, "light": PURPLE_LIGHT,
        "bullets": [
            "Embeds incoming questions and runs semantic search against the Qdrant "
            "vector store built from indexed technical PDFs.",
            "Returns a grounded answer plus the real source document(s) — erasing the "
            "20–30 minutes once lost per document search.",
            "A human reviews and signs off the synthesis for the LangSmith/audit trail.",
        ],
        "proof": "After indexing the generated PDFs into Qdrant, a real query against a "
                 "known PO returned an actual retrieved chunk and source filename "
                 "— not a 'no match' placeholder.",
    },
    {
        "title": "M5 · Hybrid Traceability", "color": TEAL, "light": TEAL_LIGHT,
        "bullets": [
            "Routes traceability/audit requests to a HYBRID cycle: Transactional agent "
            "rebuilds the SQL trail, Investigative agent searches every related document.",
            "Merges both into one Traceability Dossier — PO, supplier, reception date, "
            "FNCs, narrative and document sources — as a single HITL task.",
            "On approval, the dossier is archived as the audit record.",
        ],
        "proof": "Querying a real serial number's full history correctly routed to the "
                 "hybrid agent, returning the structured PO trail merged with "
                 "Qdrant search results in one dossier, then archived on approval.",
    },
)

# ============================================================ SLIDE 9 =====
s = add_slide()
gradient_bg(s, BLUE, NAVY_DARK)
simple_text(s, Inches(0.7), Inches(0.55), Inches(8), Inches(0.35), "THE SOLUTION", 12, RGBColor(0xCF, 0xDE, 0xF7), bold=True, spacing=150)
simple_text(s, Inches(0.65), Inches(0.95), Inches(11.5), Inches(0.9),
            "A multi-agent AI system, delivered for the factory floor", 32, WHITE, bold=True, font=FONT_HEAD)
simple_text(s, Inches(0.7), Inches(1.85), Inches(11.5), Inches(0.7),
            "Actu'AI bridges the structured/unstructured gap with specialised AI agents — all five "
            "missions automated, data kept under local control, a human always in command.",
            13, RGBColor(0xE4, 0xEC, 0xFB), line_spacing=1.15)

feats = [
    ("Specialised agents", "A semantic router dispatches work to a Transactional agent (SQL) and an Investigative agent (RAG).", BLUE),
    ("Data sovereignty", "All databases, parsing and the routing model run on air-gapped local servers.", GREEN),
    ("Human-in-the-loop", "No agent has autonomous write access. Every action is validated by a human expert.", ORANGE),
    ("Proven in production-like tests", "Exercised on a live Docker Compose stack — real DB, real vector store, real APIs.", PURPLE),
]
fw = Inches(2.78)
for i, (t, d, col) in enumerate(feats):
    fx = Inches(0.7) + i * (fw + Inches(0.1))
    cc = card(s, fx, Inches(2.85), fw, Inches(3.3), fill=WHITE)
    icon_badge(s, fx + Inches(0.25), Inches(3.1), Inches(0.55), {BLUE: BLUE_LIGHT, GREEN: GREEN_LIGHT, ORANGE: ORANGE_LIGHT, PURPLE: PURPLE_LIGHT}[col])
    simple_text(s, fx + Inches(0.25), Inches(3.85), fw - Inches(0.5), Inches(0.6), t, 13.5, NAVY, bold=True, line_spacing=1.05)
    simple_text(s, fx + Inches(0.25), Inches(4.5), fw - Inches(0.5), Inches(1.4), d, 10.8, GRAY_TXT, line_spacing=1.15)

simple_text(s, Inches(0.7), Inches(6.55), Inches(11.5), Inches(0.6),
            "The next slides walk down the delivered architecture: application services, the "
            "agents and their models, data, integration and infrastructure.", 11.5,
            RGBColor(0xCF, 0xDE, 0xF7), italic=True)

# ============================================================ SLIDE 10 ====
s = add_slide(); set_bg(s, BG)
header(s, "10", "APPLICATION SERVICES LAYER", "A modular, containerised microservices stack")

diag = card(s, Inches(0.55), Inches(1.35), Inches(5.9), Inches(5.4), fill=RGBColor(0xFB, 0xF7, 0xE7))
def svc_box(x, y, w, h, label, sub, fill, line):
    b = rect(s, x, y, w, h, fill, radius=0.08)
    b.line.color.rgb = line; b.line.width = Pt(1.25)
    simple_text(s, x, y + Inches(0.08), w, Inches(0.3), label, 11.5, NAVY, bold=True, align=PP_ALIGN.CENTER)
    simple_text(s, x, y + Inches(0.4), w, Inches(0.3), sub, 9.5, GRAY_TXT, align=PP_ALIGN.CENTER)

svc_box(Inches(2.05), Inches(1.65), Inches(2.85), Inches(0.75), "actuai_frontend", "(React)", RGBColor(0xDD, 0xEE, 0xFC), BLUE)
arrow1 = s.shapes.add_connector(1, Inches(3.475), Inches(2.4), Inches(3.475), Inches(2.85))
arrow1.line.color.rgb = GRAY_TXT; arrow1.line.width = Pt(1.5)
simple_text(s, Inches(3.75), Inches(2.45), Inches(2.0), Inches(0.4), "REST (HITL)", 8.5, GRAY_TXT, italic=True)

svc_box(Inches(1.55), Inches(2.85), Inches(3.85), Inches(0.85), "actuai_backend", "(FastAPI / LangGraph / ETL)", RGBColor(0xFC, 0xE3, 0xC9), ORANGE)

arrow2 = s.shapes.add_connector(1, Inches(2.6), Inches(3.7), Inches(2.0), Inches(4.2))
arrow2.line.color.rgb = GRAY_TXT; arrow2.line.width = Pt(1.5)
arrow3 = s.shapes.add_connector(1, Inches(4.4), Inches(3.7), Inches(5.0), Inches(4.2))
arrow3.line.color.rgb = GRAY_TXT; arrow3.line.width = Pt(1.5)

svc_box(Inches(1.05), Inches(4.2), Inches(2.85), Inches(0.95), "actuai_mock_data", "Mock SAP ERP / Generators", RGBColor(0xEB, 0xDE, 0xF7), PURPLE)
svc_box(Inches(4.05), Inches(4.2), Inches(2.4), Inches(0.95), "Local Datalake", "PostgreSQL & Qdrant", RGBColor(0xDC, 0xEF, 0xE2), GREEN)

simple_text(s, Inches(0.55), Inches(5.45), Inches(5.9), Inches(1.0),
            "Figure 1 — Microservices architecture (delivered, containerised via Docker Compose)",
            9.5, GRAY_TXT, italic=True, align=PP_ALIGN.CENTER)

rx = Inches(6.65)
items = [
    ("actuai_frontend", "React. The presentation layer and Human-in-the-Loop validation dashboard.", BLUE),
    ("actuai_backend", "FastAPI + Python. ETL pipelines, routing logic and the LangGraph orchestration.", ORANGE),
    ("actuai_mock_data", "FastAPI + SQLModel. Fake SAP ERP (BAPI) and synthetic email/PDF/Excel generators.", PURPLE),
]
for i, (t, d, col) in enumerate(items):
    cy = Inches(1.35) + i * Inches(1.25)
    cc = card(s, rx, cy, Inches(6.1), Inches(1.1))
    simple_text(s, rx + Inches(0.25), cy + Inches(0.12), Inches(5.6), Inches(0.3), t, 13, col, bold=True)
    simple_text(s, rx + Inches(0.25), cy + Inches(0.45), Inches(5.6), Inches(0.55), d, 10.5, GRAY_TXT, line_spacing=1.1)

note = rect(s, rx, Inches(5.15), Inches(6.1), Inches(1.6), NAVY, radius=0.08)
simple_text(s, rx + Inches(0.3), Inches(5.32), Inches(5.5), Inches(0.3), "Key dependency, confirmed in testing", 11, WHITE, bold=True)
simple_text(s, rx + Inches(0.3), Inches(5.65), Inches(5.5), Inches(1.0),
            "The backend is the sole orchestrator of the datalake. In the integration run, "
            "neither frontend nor mock data touched PostgreSQL/Qdrant directly — every read "
            "and write passed through the backend's REST/ETL layer as designed.",
            10, RGBColor(0xCF, 0xDE, 0xF7), line_spacing=1.15)
footer(s)

# ============================================================ SLIDE 11 ====
s = add_slide(); set_bg(s, BG)
header(s, "11", "AI AGENT LAYER", "Triggered events, orchestrated by LangGraph")
simple_text(s, Inches(0.55), Inches(1.28), Inches(11.6), Inches(0.4),
            "The orchestration cycle starts asynchronously, driven by three operational bottlenecks:",
            12, GRAY_TXT)

trigs = [
    ("Supply-chain communications", "30–60 supplier emails per employee per day on delivery status, delays and shipping.", BLUE),
    ("ERP discrepancies", "Internal alerts when SAP forecast dates fail to align with the Airbus production schedule.", ORANGE),
    ("Compliance audits", "Manual requests needing the immediate compilation of a component's traceability history.", GREEN),
]
tw = Inches(3.78)
for i, (t, d, col) in enumerate(trigs):
    tx = Inches(0.55) + i * (tw + Inches(0.08))
    cc = card(s, tx, Inches(1.8), tw, Inches(1.55))
    icon_badge(s, tx + Inches(0.25), Inches(2.0), Inches(0.5), {BLUE: BLUE_LIGHT, ORANGE: ORANGE_LIGHT, GREEN: GREEN_LIGHT}[col])
    simple_text(s, tx + Inches(0.9), Inches(2.0), tw - Inches(1.1), Inches(0.5), t, 12, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    simple_text(s, tx + Inches(0.25), Inches(2.6), tw - Inches(0.5), Inches(0.7), d, 10, GRAY_TXT, line_spacing=1.1)

flow = card(s, Inches(0.55), Inches(3.55), Inches(7.6), Inches(3.25), fill=RGBColor(0xFB, 0xF7, 0xE7))
fy = Inches(3.8)
def flow_box(x, y, w, h, label, fill, line, size=10):
    b = rect(s, x, y, w, h, fill, radius=0.1)
    b.line.color.rgb = line
    simple_text(s, x, y, w, h, label, size, NAVY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

flow_box(Inches(3.55), fy, Inches(1.6), Inches(0.55), "Global State", RGBColor(0xFD, 0xE9, 0xB8), ORANGE)
flow_box(Inches(3.4), fy + Inches(0.8), Inches(1.9), Inches(0.55), "Supervisor\nAgent", RGBColor(0xDD, 0xEE, 0xFC), BLUE, size=9.5)
flow_box(Inches(1.1), fy + Inches(1.65), Inches(2.3), Inches(0.55), "Transactional Agent\nSQL Tools", RGBColor(0xFC, 0xE3, 0xC9), ORANGE, size=9)
flow_box(Inches(5.0), fy + Inches(1.65), Inches(2.3), Inches(0.55), "Investigative Agent\nRAG Tools", RGBColor(0xEB, 0xDE, 0xF7), PURPLE, size=9)
flow_box(Inches(2.6), fy + Inches(2.45), Inches(3.0), Inches(0.55), "Human-in-the-Loop Review → Execution", RGBColor(0xDC, 0xEF, 0xE2), GREEN, size=8.5)

for (x1, y1, x2, y2) in [
    (4.35, fy/914400+0.55, 4.35, fy/914400+0.8),
]:
    pass
simple_text(s, Inches(0.85), Inches(6.55), Inches(7), Inches(0.25),
            "Figure 2 — Trigger → routing → specialist worker → human validation (as delivered)",
            9, GRAY_TXT, italic=True)

side = card(s, Inches(8.3), Inches(3.55), Inches(4.45), Inches(3.25), fill=NAVY)
simple_text(s, Inches(8.6), Inches(3.78), Inches(3.9), Inches(0.35), "Shared global state", 13.5, WHITE, bold=True)
simple_text(s, Inches(8.6), Inches(4.2), Inches(3.9), Inches(1.1),
            "On each trigger, LangGraph initialises a Global State — a shared memory "
            "dictionary holding the raw input and execution context as it moves through the graph.",
            11, RGBColor(0xCF, 0xDE, 0xF7), line_spacing=1.2)
simple_text(s, Inches(8.6), Inches(5.4), Inches(3.9), Inches(1.2),
            "Two specialist workers — plus the hybrid traceability path added for Mission 5 — "
            "now automate all five NVA-generating missions, with every result paused for "
            "human validation.", 11, RGBColor(0xCF, 0xDE, 0xF7), italic=True, line_spacing=1.2)
footer(s)

# ============================================================ SLIDE 12 ====
s = add_slide(); set_bg(s, BG)
header(s, "12", "AI AGENT LAYER · THE AGENTS", "One router, two specialist workers — plus a hybrid path")

agents = [
    ("Supervisor Agent", "Semantic Router", BLUE, BLUE_LIGHT,
     "A single entry point using a fast local LLM call for intent classification. No "
     "business logic, no database access — it only routes.",
     ["Called on every routing iteration", "Needs near-zero latency", "Runs locally, no cloud quota"]),
    ("Transactional Agent", "Structured Data · SQL Tools", ORANGE, ORANGE_LIGHT,
     "Python SQL tools query the local PostgreSQL datalake and format structured ERP "
     "metadata for high-precision tasks.",
     ["M1 · Extract email statuses, sync SAP", "M2 · Flag deltas ahead, avert AOG", "M3 · Pre-fill FNC, draft 8D follow-ups"]),
    ("Investigative Agent", "Unstructured Data · RAG Tools", PURPLE, PURPLE_LIGHT,
     "RAG tools run semantic search across the Qdrant vector database over large volumes "
     "of unstructured text and documents.",
     ["M4 · Retrieve correct, latest doc versions", "M5 · Feed the hybrid traceability dossier", "Cross-checks serials against SAP records"]),
]
aw = Inches(3.78)
for i, (t, sub, col, lcol, desc, feats_) in enumerate(agents):
    ax = Inches(0.55) + i * (aw + Inches(0.08))
    cc = card(s, ax, Inches(1.4), aw, Inches(5.35))
    icon_badge(s, ax + Inches(0.3), Inches(1.65), Inches(0.6), lcol)
    simple_text(s, ax + Inches(0.3), Inches(2.45), aw - Inches(0.6), Inches(0.4), t, 14.5, NAVY, bold=True)
    simple_text(s, ax + Inches(0.3), Inches(2.85), aw - Inches(0.6), Inches(0.3), sub, 10.5, col, bold=True)
    simple_text(s, ax + Inches(0.3), Inches(3.25), aw - Inches(0.6), Inches(1.1), desc, 10.8, GRAY_TXT, line_spacing=1.15)
    tagbox = rect(s, ax + Inches(0.3), Inches(4.45), aw - Inches(0.6), Inches(0.3), lcol, radius=0.4)
    simple_text(s, ax + Inches(0.3), Inches(4.45), aw - Inches(0.6), Inches(0.3),
                "MISSIONS HANDLED" if i > 0 else "WHY IT FITS", 9, col, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=60)
    bullets(s, ax + Inches(0.3), Inches(4.95), aw - Inches(0.6), Inches(1.7), feats_, size=10.3, gap=8, bullet_color=col)
footer(s)

# ============================================================ SLIDE 13 ====
s = add_slide(); set_bg(s, BG)
header(s, "13", "AI AGENT LAYER · VALIDATION", "No autonomous writes — confirmed in every test run")

c1 = card(s, Inches(0.55), Inches(1.4), Inches(5.6), Inches(5.35))
icon_badge(s, Inches(0.9), Inches(1.65), Inches(0.55), GREEN_LIGHT)
simple_text(s, Inches(0.9), Inches(1.65), Inches(0.55), Inches(0.55), "🛡", 16, GREEN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
simple_text(s, Inches(1.6), Inches(1.68), Inches(4.4), Inches(0.5), "Why a human stays in the loop", 14, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.9), Inches(2.45), Inches(5.0), Inches(1.6), [
    "No agent has autonomous write access to the production ERP — every drafted "
    "action is treated as a proposal, never an instruction.",
    "Across all five missions tested, zero writes occurred without an explicit "
    "approve call on the matching HITL task.",
], size=11.5, gap=10)

obs = rect(s, Inches(0.9), Inches(4.3), Inches(5.0), Inches(2.2), NAVY, radius=0.08)
simple_text(s, Inches(1.15), Inches(4.48), Inches(4.5), Inches(0.35), "Observability, confirmed", 12.5, WHITE, bold=True)
simple_text(s, Inches(1.15), Inches(4.9), Inches(4.5), Inches(1.5),
            "Every tool call, retrieval and LLM generation is traceable. In the delivered "
            "system, the append-only audit log additionally hash-chains every auth event, "
            "draft and approval — verified intact (audit_chain_intact: true) after the "
            "full test run across all five missions.",
            10.8, RGBColor(0xCF, 0xDE, 0xF7), line_spacing=1.2)

c2 = card(s, Inches(6.35), Inches(1.4), Inches(6.4), Inches(5.35))
simple_text(s, Inches(6.65), Inches(1.65), Inches(5.8), Inches(0.4), "The validation flow, as delivered", 14, NAVY, bold=True)
steps2 = [
    ("Agent drafts the action", "A worker writes a structured draft (e.g. an SAP payload or FNC) into the global state.", BLUE),
    ("Execution is paused", "LangGraph explicitly halts; nothing reaches SAP or the archive while it waits.", ORANGE),
    ("Expert reviews on the dashboard", "The draft surfaces on the React front-end for inspection, edit or rejection.", PURPLE),
    ("Validated → real execution", "Only on explicit approval does the backend push the real SAP update or archive.", GREEN),
]
for i, (t, d, col) in enumerate(steps2):
    yy = Inches(2.2) + i * Inches(1.1)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.65), yy, Inches(0.5), Inches(0.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = col; no_line(circ); circ.shadow.inherit = False
    tf = circ.text_frame; tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run_text(p, str(i+1), 13, WHITE, bold=True)
    simple_text(s, Inches(7.3), yy - Inches(0.03), Inches(5.3), Inches(0.3), t, 12, NAVY, bold=True)
    simple_text(s, Inches(7.3), yy + Inches(0.3), Inches(5.3), Inches(0.6), d, 10.2, GRAY_TXT, line_spacing=1.1)
    if i < 3:
        ln = s.shapes.add_connector(1, Inches(6.9), yy + Inches(0.5), Inches(6.9), yy + Inches(1.1))
        ln.line.color.rgb = GRAY_LINE; ln.line.width = Pt(1.5)
footer(s)

# ============================================================ SLIDE 14 ====
s = add_slide(); set_bg(s, BG)
header(s, "14", "FOUNDATION MODEL LAYER", "A hybrid LLM strategy — local speed, cloud power")
simple_text(s, Inches(0.55), Inches(1.28), Inches(11.6), Inches(0.4),
            "To balance reasoning against a tight local edge server (8 GB VRAM), each agent runs "
            "a different open-weight model. No vendor lock-in, no black box.", 12, GRAY_TXT)

models = [
    ("Llama 3.1 · 8B Instruct", "ON-PREMISE (Ollama, 4-bit)", "Supervisor Agent",
     "Called at every routing iteration — needs near-zero latency and must not burn cloud "
     "quota. Its light footprint fits the 8 GB local VRAM budget.", BLUE, BLUE_LIGHT),
    ("Mistral-Nemo · 12B", "CLOUD (NVIDIA NIM / Azure)", "Transactional Agent",
     "Chosen for its strict tool-calling compliance — reliably generates the well-formed "
     "SQL/JSON payloads the integration tests exercised.", ORANGE, ORANGE_LIGHT),
    ("Llama 3.1 70B / Mistral Large", "CLOUD (NVIDIA NIM / Azure)", "Investigative Agent",
     "A 128k-token context window ingests large RAG chunks — many supplier PDFs and email "
     "archives at once — without hallucinating.", PURPLE, PURPLE_LIGHT),
]
mw2 = Inches(3.78)
for i, (t, badge, powers, why, col, lcol) in enumerate(models):
    mx = Inches(0.55) + i * (mw2 + Inches(0.08))
    cc = card(s, mx, Inches(1.85), mw2, Inches(4.9))
    icon_badge(s, mx + Inches(0.3), Inches(2.1), Inches(0.55), lcol)
    simple_text(s, mx + Inches(0.3), Inches(2.78), mw2 - Inches(0.6), Inches(0.55), t, 14, NAVY, bold=True, line_spacing=1.0)
    tagb = rect(s, mx + Inches(0.3), Inches(3.4), mw2 - Inches(0.6), Inches(0.4), lcol, radius=0.3)
    simple_text(s, mx + Inches(0.3), Inches(3.4), mw2 - Inches(0.6), Inches(0.4), badge, 9.5, col, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    simple_text(s, mx + Inches(0.3), Inches(4.0), mw2 - Inches(0.6), Inches(0.3), "Powers: " + powers, 11, NAVY, bold=True)
    simple_text(s, mx + Inches(0.3), Inches(4.4), mw2 - Inches(0.6), Inches(0.3), "Why this model", 10, col, bold=True)
    simple_text(s, mx + Inches(0.3), Inches(4.75), mw2 - Inches(0.6), Inches(1.85), why, 10.3, GRAY_TXT, line_spacing=1.15)
footer(s)

# ============================================================ SLIDE 15 ====
s = add_slide(); set_bg(s, BG)
header(s, "15", "FOUNDATION MODEL LAYER · BALANCE", "What the hybrid design buys — and its trade-offs")

c1 = card(s, Inches(0.55), Inches(1.4), Inches(5.6), Inches(5.35), fill=GREEN_LIGHT, line=GREEN_LIGHT)
icon_badge(s, Inches(0.9), Inches(1.65), Inches(0.5), WHITE)
simple_text(s, Inches(0.9), Inches(1.65), Inches(0.5), Inches(0.5), "✓", 15, GREEN, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
simple_text(s, Inches(1.55), Inches(1.68), Inches(4.4), Inches(0.5), "Main capabilities", 14.5, RGBColor(0x14,0x5A,0x2C), bold=True, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.9), Inches(2.4), Inches(5.0), Inches(4.1), [
    ("Specialised division of labour — ", "fast, deterministic routing locally; heavy synthesis and extraction in the cloud, confirmed live during testing."),
    ("Open-weight throughout — ", "every selected model is open-weight: no vendor lock-in, no proprietary black box."),
    ("Latency where it matters — ", "the router never pays network cost, keeping classification instant on local hardware."),
], size=11.5, gap=12, bullet_color=GREEN, color=RGBColor(0x1F,0x4A,0x2E))

c2 = card(s, Inches(6.35), Inches(1.4), Inches(6.4), Inches(5.35), fill=ORANGE_LIGHT, line=ORANGE_LIGHT)
icon_badge(s, Inches(6.7), Inches(1.65), Inches(0.5), WHITE)
simple_text(s, Inches(6.7), Inches(1.65), Inches(0.5), Inches(0.5), "!", 16, ORANGE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
simple_text(s, Inches(7.35), Inches(1.68), Inches(5.2), Inches(0.5), "Limitations to manage", 14.5, RGBColor(0x8A,0x49,0x0E), bold=True, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(6.7), Inches(2.4), Inches(5.8), Inches(2.5), [
    ("Local hardware bottleneck — ", "the local Llama 3.1 8B is 4-bit quantized, slightly degrading reasoning vs. uncompressed."),
    ("Cloud dependency — ", "an NVIDIA NIM/Azure outage temporarily paralyses the Transactional and Investigative workflows."),
], size=11.5, gap=12, bullet_color=ORANGE, color=RGBColor(0x6B,0x3D,0x12))
mit = rect(s, Inches(6.7), Inches(4.95), Inches(5.7), Inches(1.6), WHITE, radius=0.08)
simple_text(s, Inches(6.9), Inches(5.1), Inches(5.3), Inches(0.3), "Mitigation, observed in practice", 11, ORANGE, bold=True)
simple_text(s, Inches(6.9), Inches(5.45), Inches(5.3), Inches(1.0),
            "The human-in-the-loop checkpoint and audit logging kept every test run "
            "observable and recoverable — no silent failures were seen end-to-end.",
            10.3, GRAY_TXT, line_spacing=1.15)
footer(s)

# ============================================================ SLIDE 16 ====
s = add_slide(); set_bg(s, BG)
header(s, "16", "DATA LAYER", "Unifying structured & unstructured sources")

blocks = [
    ("Structured — source of truth", "SAP ERP holds forecast dates, part references, supplier names, PO numbers and "
     "serial numbers. Agents read this live to replace manual entry and flag discrepancies.", BLUE, BLUE_LIGHT),
    ("Unstructured — the daily flow", "30–60 emails/day on delivery status, archives of past non-conformities, "
     "regulatory PDFs, and legacy signed paper records.", ORANGE, ORANGE_LIGHT),
    ("Local datalake — dual database", "PostgreSQL + SQLModel store operational logs and SAP metadata; Qdrant stores "
     "embeddings for RAG. Both confirmed populated and queried during testing.", PURPLE, PURPLE_LIGHT),
]
bw = Inches(3.78)
for i, (t, d, col, lcol) in enumerate(blocks):
    bx = Inches(0.55) + i * (bw + Inches(0.08))
    cc = card(s, bx, Inches(1.4), bw, Inches(2.1))
    icon_badge(s, bx + Inches(0.25), Inches(1.63), Inches(0.5), lcol)
    simple_text(s, bx + Inches(0.9), Inches(1.63), bw - Inches(1.1), Inches(0.5), t, 12, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    simple_text(s, bx + Inches(0.25), Inches(2.25), bw - Inches(0.5), Inches(1.1), d, 10.3, GRAY_TXT, line_spacing=1.15)

diag2 = card(s, Inches(0.55), Inches(3.7), Inches(11.6), Inches(3.05), fill=WHITE)
labels2 = [
    ("Physical Docs\nScanned & Archived", GRAY_TXT), ("Network Drives\nPDFs, 8D, Reports", GRAY_TXT),
    ("MS Exchange\n30-60 Emails/day", GRAY_TXT), ("Excel Files\nShared Dashboards", GRAY_TXT),
    ("SAP ERP\nPP, MM, QM", GRAY_TXT),
]
for i, (lab, col) in enumerate(labels2):
    ly = Inches(3.9) + i * Inches(0.5)
    bb = rect(s, Inches(0.85), ly, Inches(1.9), Inches(0.42), RGBColor(0xEE, 0xF1, 0xF6), radius=0.2)
    simple_text(s, Inches(0.85), ly, Inches(1.9), Inches(0.42), lab.split("\n")[0], 8, GRAY_TXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

etl = rect(s, Inches(3.1), Inches(3.85), Inches(2.6), Inches(2.65), RGBColor(0xFB, 0xF7, 0xE7), radius=0.06)
simple_text(s, Inches(3.1), Inches(3.95), Inches(2.6), Inches(0.3), "ETL Pipeline", 10.5, ORANGE, bold=True, align=PP_ALIGN.CENTER)
simple_text(s, Inches(3.25), Inches(4.3), Inches(2.3), Inches(2.0),
            "Connectors & extractors\n(SAP API, IMAP, OCR)\n\nCleaning & validation\n(strict typing)\n\n"
            "Embedding model\n(semantic vectorization)", 8.7, GRAY_TXT, align=PP_ALIGN.CENTER, line_spacing=1.1)

db1 = rect(s, Inches(6.3), Inches(4.0), Inches(2.15), Inches(0.95), RGBColor(0xDC, 0xEF, 0xE2), radius=0.1)
simple_text(s, Inches(6.3), Inches(4.0), Inches(2.15), Inches(0.95), "PostgreSQL\nStructured Metadata", 9.5, RGBColor(0x14,0x5A,0x2C), bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
db2 = rect(s, Inches(6.3), Inches(5.15), Inches(2.15), Inches(0.95), RGBColor(0xEB, 0xDE, 0xF7), radius=0.1)
simple_text(s, Inches(6.3), Inches(5.15), Inches(2.15), Inches(0.95), "Qdrant\nVectors & Documents", 9.5, RGBColor(0x4A,0x32,0x8A), bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)

ag1 = rect(s, Inches(9.2), Inches(4.0), Inches(2.5), Inches(0.95), RGBColor(0xFC, 0xE3, 0xC9), radius=0.1)
simple_text(s, Inches(9.2), Inches(4.0), Inches(2.5), Inches(0.95), "Transactional Agent\nExact SQL Queries", 9.5, RGBColor(0x8A,0x49,0x0E), bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
ag2 = rect(s, Inches(9.2), Inches(5.15), Inches(2.5), Inches(0.95), RGBColor(0xDD, 0xEE, 0xFC), radius=0.1)
simple_text(s, Inches(9.2), Inches(5.15), Inches(2.5), Inches(0.95), "Investigative Agent\nRAG / Semantic Search", 9.5, RGBColor(0x1A,0x3D,0x7A), bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)

simple_text(s, Inches(0.55), Inches(6.5), Inches(11.6), Inches(0.2),
            "Figure 3 — Dispersed sources → ETL → on-premise datalake → agents (delivered pipeline)",
            9, GRAY_TXT, italic=True, align=PP_ALIGN.CENTER)
footer(s)

# ============================================================ SLIDE 17 ====
s = add_slide(); set_bg(s, BG)
header(s, "17", "INTEGRATION LAYER", "How the pieces talk — confirmed end-to-end")

cols17 = [
    ("Internal APIs & connectors", BLUE, BLUE_LIGHT, [
        ("Frontend ↔ backend — ", "the React dashboard talks to the orchestration logic over secure REST APIs, verified live through the dev-server proxy."),
        ("Database connectors — ", "SQLAlchemy/SQLModel reads/writes PostgreSQL; the Qdrant client indexes and retrieves vectors."),
    ]),
    ("External systems", ORANGE, ORANGE_LIGHT, [
        ("SAP BAPI — ", "the ETL pulls metadata via HTTP GET; validated updates push back with HTTP PUT/POST — confirmed against the mock SAP."),
        ("MS Exchange webhooks — ", "supplier emails fire HTTP POST to the ingestion router, launching the LangGraph cycle automatically."),
    ]),
    ("Cloud AI integration", PURPLE, PURPLE_LIGHT, [
        ("Secure enterprise APIs — ", "heavy models are reached on NVIDIA NIM / Azure through stateless, API-key-protected calls."),
        ("Zero-data-retention — ", "strict no-retention policies keep sensitive aerospace data out of training pipelines."),
    ]),
]
cw = Inches(3.78)
for i, (t, col, lcol, items) in enumerate(cols17):
    cx = Inches(0.55) + i * (cw + Inches(0.08))
    cc = card(s, cx, Inches(1.4), cw, Inches(5.35))
    icon_badge(s, cx + Inches(0.3), Inches(1.65), Inches(0.55), lcol)
    simple_text(s, cx + Inches(0.3), Inches(2.35), cw - Inches(0.6), Inches(0.45), t, 13.5, NAVY, bold=True, line_spacing=1.0)
    yy = 2.9
    for sub, desc in items:
        bb = rect(s, cx + Inches(0.3), Inches(yy), cw - Inches(0.6), Inches(1.55), lcol, radius=0.08)
        simple_text(s, cx + Inches(0.5), Inches(yy + 0.12), cw - Inches(1.0), Inches(0.3), sub.strip(" —"), 10.5, col, bold=True)
        simple_text(s, cx + Inches(0.5), Inches(yy + 0.45), cw - Inches(1.0), Inches(1.0), desc, 9.8, GRAY_TXT, line_spacing=1.12)
        yy += 1.75
footer(s)

# ============================================================ SLIDE 18 ====
s = add_slide(); set_bg(s, BG)
header(s, "18", "SECURITY, GOVERNANCE & INFRASTRUCTURE", "Sovereignty by design, ops-ready")

c1 = card(s, Inches(0.55), Inches(1.4), Inches(5.7), Inches(5.35))
simple_text(s, Inches(0.85), Inches(1.63), Inches(5.1), Inches(0.4), "Security & governance", 14, NAVY, bold=True)
bullets(s, Inches(0.85), Inches(2.1), Inches(5.1), Inches(4.5), [
    ("Authentication — ", "JWT-based identity with role + clearance claims; demo roles (engineer, buyer, operator_admin, auditor) seeded and exercised."),
    ("Authorization (RBAC+ABAC) — ", "only engineer/buyer/operator_admin roles can approve a write to SAP — confirmed by testing role-restricted endpoints."),
    ("Immutable audit log — ", "every draft, approval and rejection is hash-chained; integrity verified via the audit-verify endpoint after the full test run."),
    ("Compliance posture — ", "EN9100/AS9100-aligned traceability, human-always-signs stance for sensitive write-backs, in line with ITAR/EASA principles from the design."),
], size=10.8, gap=10)

c2 = card(s, Inches(6.4), Inches(1.4), Inches(6.35), Inches(5.35))
simple_text(s, Inches(6.7), Inches(1.63), Inches(5.7), Inches(0.4), "Infrastructure, as deployed", 14, NAVY, bold=True)
bullets(s, Inches(6.7), Inches(2.1), Inches(5.7), Inches(4.5), [
    ("Hybrid design — ", "on-premise edge zone (DB, business logic, router model) reaching a managed cloud zone for heavy models only."),
    ("Containerised compute — ", "Docker Compose orchestrates frontend, backend, mock-data, PostgreSQL and Qdrant — all started, healthy and verified together."),
    ("Isolated networking — ", "PostgreSQL/Qdrant expose no host ports; only the backend container reaches them."),
    ("Resilience — ", "Docker restart-unless-stopped policies; cron-job dumps for PostgreSQL and Qdrant snapshots."),
], size=10.8, gap=10)
footer(s)

# ============================================================ SLIDE 19 ====
s = add_slide(); set_bg(s, BG)
header(s, "19", "PROOF OF DELIVERY", "Validated end-to-end on a live Docker Compose stack")
simple_text(s, Inches(0.55), Inches(1.28), Inches(11.6), Inches(0.45),
            "Beyond the 11-test automated backend suite, every mission was exercised against the "
            "real stack — actual HTTP calls, real PostgreSQL/Qdrant, the simulated SAP ERP.",
            12, GRAY_TXT, line_spacing=1.15)

rows = [
    ("UC1", "Supply Chain Monitoring", "Delay email → SAP_UPDATE draft → approved → SAP delivery date confirmed updated.", GREEN),
    ("UC2", "AOG Detection", "Parallel AOG_ALERT task created automatically → approved → escalation email confirmed in outbox.", GREEN),
    ("UC3", "Quality / FNC", "\"Créer FNC...\" → pre-filled draft → approved → Quality Notification confirmed created in SAP.", GREEN),
    ("UC4", "Documentation RAG", "PDFs generated & indexed in Qdrant → real query → grounded answer + real source retrieved.", GREEN),
    ("UC5", "Hybrid Traceability", "Serial-number audit request → hybrid SQL+RAG dossier compiled → approved → archived.", GREEN),
]
ty = 1.95
for i, (tag, title, desc, col) in enumerate(rows):
    cc = card(s, Inches(0.55), Inches(ty), Inches(11.2), Inches(0.85))
    tagb = rect(s, Inches(0.8), Inches(ty + 0.18), Inches(0.85), Inches(0.5), GREEN_LIGHT, radius=0.3)
    simple_text(s, Inches(0.8), Inches(ty + 0.18), Inches(0.85), Inches(0.5), tag, 12, GREEN, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    simple_text(s, Inches(1.85), Inches(ty + 0.1), Inches(2.5), Inches(0.65), title, 12, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    simple_text(s, Inches(4.45), Inches(ty + 0.1), Inches(7.0), Inches(0.65), desc, 10.5, GRAY_TXT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    simple_text(s, Inches(0.55), Inches(ty), Inches(0.0), Inches(0.0), "", 1, WHITE)
    ck = simple_text(s, Inches(11.3), Inches(ty + 0.18), Inches(0.4), Inches(0.5), "✓", 18, GREEN, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ty += 0.95

bottom = rect(s, Inches(0.55), Inches(ty + 0.05), Inches(11.2), Inches(0.55), NAVY, radius=0.1)
simple_text(s, Inches(0.55), Inches(ty + 0.05), Inches(11.2), Inches(0.55),
            "A real bug was found and fixed during this pass (SAP-mock date typing) and the audit-log hash chain was confirmed intact afterward.",
            10.5, WHITE, italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ============================================================ SLIDE 20 ====
s = add_slide()
gradient_bg(s, BLUE, NAVY_DARK)
simple_text(s, Inches(0.7), Inches(0.55), Inches(8), Inches(0.35), "CONCLUSION", 12, RGBColor(0xCF, 0xDE, 0xF7), bold=True, spacing=150)
simple_text(s, Inches(0.65), Inches(0.95), Inches(11.8), Inches(0.9),
            "From hours of admin to instant, governed action", 30, WHITE, bold=True, font=FONT_HEAD)
simple_text(s, Inches(0.7), Inches(1.8), Inches(11.5), Inches(0.7),
            "Actu'AI has turned the actuation service's structured/unstructured gap into an "
            "automated, auditable pipeline — delivered, integration-tested, and ready for review.",
            13, RGBColor(0xE4, 0xEC, 0xFB), line_spacing=1.15)

concl = [
    ("NVA eliminated", "Reading, retyping, chasing and searching are replaced by automatic drafts across all five missions.", BLUE),
    ("Proactive, not reactive", "Deliveries are cross-checked continuously, flagging delays ahead of time and preventing AOG crises.", ORANGE),
    ("Sovereign & compliant", "Air-gapped data, zero-retention cloud calls, and an intact hash-chained audit log, by design.", GREEN),
    ("Human in command", "Open-weight models and a mandatory validation step kept experts in control of every write, every time.", PURPLE),
]
cw2 = Inches(2.78)
for i, (t, d, col) in enumerate(concl):
    cx = Inches(0.7) + i * (cw2 + Inches(0.1))
    cc = card(s, cx, Inches(2.75), cw2, Inches(3.1), fill=WHITE)
    icon_badge(s, cx + Inches(0.25), Inches(3.0), Inches(0.55), {BLUE: BLUE_LIGHT, ORANGE: ORANGE_LIGHT, GREEN: GREEN_LIGHT, PURPLE: PURPLE_LIGHT}[col])
    simple_text(s, cx + Inches(0.25), Inches(3.7), cw2 - Inches(0.5), Inches(0.55), t, 13, NAVY, bold=True, line_spacing=1.0)
    simple_text(s, cx + Inches(0.25), Inches(4.3), cw2 - Inches(0.5), Inches(1.4), d, 10.5, GRAY_TXT, line_spacing=1.15)

bar = rect(s, Inches(0.7), Inches(6.2), Inches(11.93), Inches(0.85), NAVY_DARK, radius=0.08)
tb, tf = textbox(s, Inches(1.0), Inches(6.35), Inches(11.4), Inches(0.6))
p = tf.paragraphs[0]
run_text(p, "Actu'AI · ECE Paris   ", 11, RGBColor(0xCF, 0xDE, 0xF7), bold=True)
run_text(p, "Jarfino HOUNGBADJI · SOW Achta Demba · Fatoumata IBRAHIM BELKO   ", 11, WHITE)
run_text(p, "github.com/jaja07/ActuAI", 11, RGBColor(0xCF, 0xDE, 0xF7), bold=True)

import os
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ActuAI_Presentation_Final.pptx")
prs.save(out_path)
print("Saved:", out_path)
